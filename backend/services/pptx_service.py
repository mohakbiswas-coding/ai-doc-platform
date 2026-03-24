# backend/services/pptx_service.py

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import io


def build_pptx(project_title: str, topic: str, sections: list) -> bytes:
    """
    Build a .pptx file from project sections.
    Returns file as bytes.
    """
    prs = Presentation()

    # Slide dimensions (widescreen 16:9)
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # ── Title slide ───────────────────────────────────────────────────────────
    title_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(title_layout)

    slide.shapes.title.text = project_title
    slide.placeholders[1].text = topic

    # ── Content slides ────────────────────────────────────────────────────────
    content_layout = prs.slide_layouts[1]  # Title and Content layout

    for section in sections:
        slide = prs.slides.add_slide(content_layout)

        # Set slide title
        slide.shapes.title.text = section["title"]

        # Set slide content
        content = section.get("content", "")
        tf = slide.placeholders[1].text_frame
        tf.clear()

        if content:
            lines = [line.strip() for line in content.split("\n") if line.strip()]
            for i, line in enumerate(lines):
                # Remove bullet character if AI included it
                clean_line = line.lstrip("•-* ").strip()
                if i == 0:
                    tf.paragraphs[0].text = clean_line
                else:
                    p = tf.add_paragraph()
                    p.text = clean_line
                    p.level = 0

    # ── Save to bytes ─────────────────────────────────────────────────────────
    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer.getvalue()
